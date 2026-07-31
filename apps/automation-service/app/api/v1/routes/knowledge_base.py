"""Bounded knowledge-file parsing and AI rule extraction routes.

Supported formats are UTF-8/GB18030 text, CSV, and the OOXML .xlsx/.pptx
containers handled by the pinned parsers. Legacy binary .xls/.ppt files are
rejected explicitly instead of advertising a capability the runtime lacks.
"""
import asyncio
import io
import csv
import logging
import stat
import zipfile
from pathlib import PurePosixPath
from typing import Any, Dict, Optional

from fastapi import APIRouter, Depends, File, Form, Header, UploadFile
from sqlalchemy.ext.asyncio import AsyncSession

from ....core.database import get_db
from ....core.http_failures import log_route_failure, safe_route_failure
from ....core.response import ResultObject
from ....services.ai_provider import generate_text
from ....services.ai_billing import (
    AiBillingError,
    AiBillingPaymentRequired,
    charge_text_usage,
    estimate_text_tokens,
    precheck_ai_usage,
)
from ....services.upload_governance import UploadGovernanceError, govern_transient_upload
from .internal import verify_internal_token
from ..deps import get_current_user

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/knowledge-base", tags=["knowledgeBase"])

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
MAX_OOXML_MEMBERS = 512
MAX_OOXML_UNCOMPRESSED_BYTES = 50 * 1024 * 1024
MAX_OOXML_COMPRESSION_RATIO = 100
MAX_OOXML_MEMBER_BYTES = 20 * 1024 * 1024
MAX_PARSED_CHARACTERS = 100_000
MAX_WORKSHEETS = 50
MAX_SPREADSHEET_ROWS = 10_000
MAX_SPREADSHEET_CELLS = 100_000
MAX_SLIDES = 200
MAX_PRESENTATION_SHAPES = 5_000
MAX_PRESENTATION_TABLE_CELLS = 50_000
PARSE_TIMEOUT_SECONDS = 15.0
_PARSE_CONCURRENCY = 2
_parse_semaphore = asyncio.Semaphore(_PARSE_CONCURRENCY)

# Legacy binary .xls/.ppt files are deliberately not advertised: the pinned
# parsers support OOXML (.xlsx/.pptx), not the older OLE container formats.
# .docx uses python-docx; .pdf uses pdfplumber.
ALLOWED_EXTENSIONS = {".md", ".txt", ".pptx", ".xlsx", ".csv", ".docx", ".pdf"}

_ALLOWED_MEDIA_TYPES = {
    ".md": {"application/octet-stream", "text/markdown", "text/plain", "text/x-markdown"},
    ".txt": {"application/octet-stream", "text/plain"},
    ".csv": {
        "application/csv",
        "application/octet-stream",
        "application/vnd.ms-excel",
        "text/csv",
        "text/plain",
    },
    ".xlsx": {
        "application/octet-stream",
        "application/zip",
        "application/x-zip-compressed",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    },
    ".pptx": {
        "application/octet-stream",
        "application/zip",
        "application/x-zip-compressed",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    },
    ".docx": {
        "application/octet-stream",
        "application/zip",
        "application/x-zip-compressed",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    },
    ".pdf": {
        "application/octet-stream",
        "application/pdf",
    },
}

# 扩展名 → 文件类型展示名（用于前端展示）
EXTENSION_TYPE_LABEL = {
    ".md": "Markdown",
    ".txt": "文本",
    ".csv": "CSV 表格",
    ".xlsx": "Excel 表格",
    ".pptx": "PPT 幻灯片",
    ".docx": "Word 文档",
    ".pdf": "PDF 文档",
}


class KnowledgeUploadError(ValueError):
    def __init__(self, public_message: str, status_code: int = 422):
        self.public_message = public_message
        self.status_code = status_code
        super().__init__(public_message)

EXTRACT_PROMPT_TEMPLATE = """你是客服规则提取助手。请从以下文件内容中提取所有可作为 AI 客服回复规则的信息，输出为结构化 Markdown 文本。

要求：
1. 按类别分组，使用二级标题（如 ## 售后政策 / ## 发货说明 / ## 商品 FAQ / ## 退换货规则 / ## 价格优惠 / ## 规格参数）
2. 每条规则用 "- " 开头，包含：触发场景、回复要点、注意事项
3. 只输出与客服回复相关的内容，忽略文件中的导航、版权、广告等无关信息
4. 保持原文事实，不要编造规则
5. 如果文件内容与客服无关，返回空字符串

文件内容：
{file_content}
"""


# Q&A 提取 prompt（用于"新建知识库"弹窗的文件上传模式，返回 JSON 数组）
# 根据文件类型（file_type）选择专用模板：
#   - auto：自动检测（先识别后用对应模板）
#   - chat_records：聊天记录文件
#   - product_docs：商品资料/说明书
#   - company_docs：公司资料/政策文档
#   - general：通用（fallback）

# 通用 fallback 提取 prompt（保持原有行为）
QA_EXTRACT_PROMPT_TEMPLATE = """你是客服知识库提取助手。请从以下文件内容中提取所有可作为客服问答对（Q&A）的高价值知识，输出为 JSON 数组。

要求：
1. 仅提取可作为"买家提问 → 卖家回复"的客服问答对，跳过导航、版权、广告、目录等无关内容
2. 每条 Q&A 包含：
   - question：买家的典型提问（简洁清晰，不超过 100 字；如原文是陈述句请改写为疑问句）
   - answer：卖家的标准回复（完整、可执行，不超过 5000 字；保留原文事实，不要编造）
   - category：分类名称（从下列 13 个一级分类中选择一个，严禁自创）：
     * 交易通用 / 服饰鞋包 / 数码家电 / 美妆个护 / 家居生活 / 母婴用品
     * 运动户外 / 图书教材 / 艺术品收藏 / 宠物用品 / 汽车用品 / 手工DIY / 虚拟货源
   - tags：3-5 个标签，逗号分隔（如：正版,售后,退款）
   - source_summary：一句话说明该 Q&A 的来源段落或主题
3. 对敏感信息脱敏：手机号 → [手机号]，微信号/QQ号 → [联系方式]，收货地址 → [地址]，真实姓名 → [姓名]
4. 如果文件内容与客服无关，返回空数组 []
5. 最多提取 30 条 Q&A，优先提取高价值、高频问题
6. 输出严格的 JSON 数组格式，无其他文字、无 markdown 代码块包装

输出格式示例：
[{{"question":"这本书是正版吗？","answer":"是的，本店所售图书均为正版，支持专柜验货。","category":"图书教材","tags":"正版,验货,图书","source_summary":"正版承诺说明"}}]

文件内容：
{file_content}
"""


# 聊天记录专用提取 prompt
# 适用场景：导出的微信/QQ/闲鱼聊天记录、客服对话日志
# 核心策略：识别买卖双方角色，抽取真实问答对，保留销售技巧
CHAT_RECORDS_QA_EXTRACT_PROMPT_TEMPLATE = """你是客服知识库提取助手。以下是导出的聊天记录（可能来自微信/QQ/闲鱼等平台）。请从中提取高价值的客服问答对（Q&A），输出为 JSON 数组。

聊天记录特点：
- 可能包含时间戳、发送者昵称、消息内容
- 买卖双方角色可能交替发言，需识别"买家提问 → 卖家回复"的真实对话
- 同一主题可能跨多条消息展开（多轮交互）

提取策略：
1. 识别买卖双方角色（通常卖家 = 商家/客服/店主，买家 = 询问商品的人）
2. 将分散的多轮对话合并为完整 Q&A：
   - 买家的多段提问合并为一个完整问题
   - 卖家的多段回复合并为一个完整回答
   - 跳过纯问候、表情、砍价无果、确认收货等无价值对话
3. 优先提取以下高价值内容：
   - 商品规格/参数咨询（尺寸、材质、容量、型号等）
   - 价格/优惠/活动咨询
   - 售后政策（退换货、保修、发票）
   - 物流发货咨询（发货时间、快递、到货）
   - 商品真伪/正版验证咨询
   - 使用方法/操作指导
4. 每条 Q&A 包含：
   - question：买家提问（改写为典型疑问句，不超过 100 字）
   - answer：卖家回复（保留原文要点，可整合多段，不超过 5000 字）
   - category：13 个一级分类之一（交易通用 / 服饰鞋包 / 数码家电 / 美妆个护 / 家居生活 / 母婴用品 / 运动户外 / 图书教材 / 艺术品收藏 / 宠物用品 / 汽车用品 / 手工DIY / 虚拟货源）
   - tags：3-5 个标签
   - source_summary：来源聊天片段的简短描述（如"买家询问发货时间"）
5. 对敏感信息脱敏：手机号 → [手机号]，微信号/QQ号 → [联系方式]，收货地址 → [地址]，真实姓名 → [姓名]
6. 最多提取 30 条 Q&A
7. 输出严格的 JSON 数组格式，无其他文字、无 markdown 代码块包装

输出格式示例：
[{{"question":"这本书什么时候发货？","answer":"工作日当天 17 点前下单当天发货，节假日顺延，默认发中通快递。","category":"图书教材","tags":"发货,快递,时效","source_summary":"买家询问发货时间"}}]

聊天记录内容：
{file_content}
"""


# 商品资料专用提取 prompt
# 适用场景：商品说明书、产品手册、规格表、商品详情页文案
# 核心策略：将陈述性商品信息改写为"买家会问什么 → 标准回复"
PRODUCT_DOCS_QA_EXTRACT_PROMPT_TEMPLATE = """你是客服知识库提取助手。以下是商品资料/产品说明书/规格表。请将商品信息改写为客服问答对（Q&A），输出为 JSON 数组。

商品资料特点：
- 多为陈述性描述（如"本产品支持 5V/2A 充电"、"材质：304 不锈钢"）
- 包含规格、参数、功能、使用方法、注意事项
- 缺少对话形式，需要从买家视角"反向生成"提问

提取策略：
1. 从商品资料中识别可被买家询问的信息点：
   - 规格参数（尺寸、重量、容量、颜色、型号）
   - 材质成分（面料、材质、配料表、成分）
   - 功能特性（支持什么、不支持什么、兼容性）
   - 使用方法（如何操作、如何安装、注意事项）
   - 包装清单（含哪些配件、是否含电池）
   - 产地/品牌/认证信息
   - 价格/促销/库存（如资料中提及）
2. 将每个信息点改写为"买家典型提问 → 卖家标准回复"：
   - question：模拟买家会怎么问（口语化、自然）
   - answer：基于资料的准确回复（保留原文事实，不要编造未提及的信息）
3. 对于表格/规格表，逐行转换为 Q&A
4. 每条 Q&A 包含：
   - question：买家典型提问（不超过 100 字）
   - answer：基于资料的回复（不超过 5000 字）
   - category：13 个一级分类之一
   - tags：3-5 个标签（含商品相关关键词）
   - source_summary：来源段落或规格项
5. 如果资料与客服无关或过于简短，返回空数组 []
6. 最多提取 30 条 Q&A，覆盖资料中所有可询问的信息点
7. 输出严格的 JSON 数组格式，无其他文字、无 markdown 代码块包装

输出格式示例：
[{{"question":"这个保温杯容量多大？","answer":"本款保温杯容量为 500ml，可满足日常饮水需求。","category":"家居生活","tags":"保温杯,容量,500ml","source_summary":"规格表-容量项"}}]

商品资料内容：
{file_content}
"""


# 公司资料专用提取 prompt
# 适用场景：公司简介、售后政策、退换货规则、发货说明、服务流程文档
# 核心策略：从政策/流程文档中提取可复用的客服规则问答
COMPANY_DOCS_QA_EXTRACT_PROMPT_TEMPLATE = """你是客服知识库提取助手。以下是公司资料/政策文档/服务流程说明。请从中提取可作为客服回复依据的问答对（Q&A），输出为 JSON 数组。

公司资料特点：
- 多为制度性、规范性内容（退换货政策、发货流程、售后规则）
- 语言正式，需转换为买家易懂的口语化回复
- 同一政策可能适用于多个咨询场景

提取策略：
1. 识别文档中的可复用客服规则：
   - 售后政策（退换货条件、时效、流程）
   - 发货说明（发货时间、快递、运费、区域限制）
   - 价格优惠（折扣、满减、会员价）
   - 服务承诺（质量保证、正品保证、发票政策）
   - 投诉与维权流程
   - 会员/积分/优惠规则
   - 营业时间/客服时间
2. 将每条规则改写为"买家咨询场景 → 标准回复"：
   - question：模拟买家会怎么咨询（如"支持 7 天无理由退货吗？"）
   - answer：基于政策文档的准确回复（保留原文要点，语言口语化但事实严谨）
3. 同一政策可拆分为多个 Q&A（如退换货政策 → 退货条件、退货流程、退款时效三个 Q&A）
4. 每条 Q&A 包含：
   - question：买家咨询场景（不超过 100 字）
   - answer：基于政策的标准回复（不超过 5000 字）
   - category：13 个一级分类之一（多数为"交易通用"）
   - tags：3-5 个标签（如：退货,7天无理由,售后）
   - source_summary：来源政策段落
5. 如果文档与客服无关，返回空数组 []
6. 最多提取 30 条 Q&A，覆盖所有可复用政策
7. 输出严格的 JSON 数组格式，无其他文字、无 markdown 代码块包装

输出格式示例：
[{{"question":"支持 7 天无理由退货吗？","answer":"支持。商品签收后 7 天内，保持商品全新未使用且包装完整，可申请无理由退货。","category":"交易通用","tags":"退货,7天无理由,售后","source_summary":"售后政策-退货条件"}}]

公司资料内容：
{file_content}
"""


# 文件类型自动检测 prompt（轻量级，仅返回类型标识）
FILE_TYPE_DETECT_PROMPT_TEMPLATE = """请分析以下文件内容片段，判断它属于哪种类型，只返回类型标识符（一个单词），不要其他文字。

类型定义：
- chat_records：聊天记录/对话日志（特征：包含时间戳、发送者昵称、来回对话、买卖双方交流）
- product_docs：商品资料/产品说明书/规格表（特征：商品规格、参数、功能描述、使用说明、包装清单）
- company_docs：公司资料/政策文档（特征：公司简介、售后政策、退换货规则、发货说明、服务流程）
- general：无法明确归类（通用文本）

文件内容片段（最多 2000 字）：
{file_content}
"""


# 会话 Q&A 提取 prompt（用于"新建知识库"弹窗的会话聊天提取模式）
# 优化点：
#   1. 按会话主题分类（售前咨询/规格询问/价格优惠/售后问题/物流发货/商品验证/使用指导）
#   2. 多轮问答合并：同一主题的多轮对话合并为完整 Q&A，避免碎片化
#   3. 跨会话去重：相似问题保留信息最完整的版本
CONVERSATION_QA_EXTRACT_PROMPT_TEMPLATE = """你是客服知识库提取助手。以下是若干真实的买家-卖家会话记录。请从中提取高价值的客服问答对（Q&A），输出为 JSON 数组。

【核心提取策略】

1. **会话主题识别**：先识别每个会话的核心主题，按主题归类提取：
   - 售前咨询（商品是否有货、是否能优惠、是否支持某功能）
   - 规格询问（尺寸、材质、容量、型号、颜色、版本）
   - 价格优惠（折扣、满减、包邮、议价结果）
   - 售后问题（退换货、保修、维修、补发）
   - 物流发货（发货时间、快递、到货时效、运费）
   - 商品验证（正版/真伪、授权、防伪）
   - 使用指导（操作方法、安装、注意事项）
   - 投诉维权（差评、纠纷、客服介入）

2. **多轮问答合并**（重要）：同一主题下买家的多段提问 + 卖家的多段回复必须合并为一条完整 Q&A：
   - 错误做法：把"这个有红色吗？"和"红色什么时候发货？"拆成两条
   - 正确做法：合并为 "这个有红色吗？红色款什么时候发货？" → "有红色，红色款工作日 17 点前下单当天发货"
   - 同一主题内多轮往复合并为一条 Q&A，保留完整上下文

3. **跨会话去重**：不同会话中相似的问题，只保留信息最完整、回复最详尽的版本，避免知识库冗余。

【输出要求】

1. 仅提取能体现真实销售技巧、产品知识、问题解决能力的高价值 Q&A，跳过纯闲聊、问候、表情、无意义对话
2. 每条 Q&A 包含：
   - question：买家的典型提问（不超过 100 字；如原文是陈述句请改写为疑问句；多轮提问合并为一个完整问题）
   - answer：卖家的优质回复（不超过 5000 字；多轮回复合并为一个完整回答；保留原文事实，可适度整理语言）
   - category：分类名称（从下列 13 个一级分类中选择一个，严禁自创）：
     * 交易通用 / 服饰鞋包 / 数码家电 / 美妆个护 / 家居生活 / 母婴用品
     * 运动户外 / 图书教材 / 艺术品收藏 / 宠物用品 / 汽车用品 / 手工DIY / 虚拟货源
   - tags：3-5 个标签，逗号分隔（建议包含主题标签，如：售前,规格,发货）
   - source_summary：来源会话的买家昵称或商品标题
   - source_conversation：来源会话标识（从输入的"对话 N"中取 N）
3. 对敏感信息脱敏：手机号 → [手机号]，微信号/QQ号 → [联系方式]，收货地址 → [地址]，真实姓名 → [姓名]
4. 如果所有对话都无价值，返回空数组 []
5. 最多提取 30 条 Q&A，优先提取高频问题、产品知识、售后政策
6. 输出严格的 JSON 数组格式，无其他文字、无 markdown 代码块包装

输出格式示例：
[{{"question":"这本书是正版吗？支持验货吗？","answer":"是的，本店所售图书均为正版，支持专柜验货。如有质量问题可七天无理由退换。","category":"图书教材","tags":"正版,验货,售后","source_summary":"张三 - 售后咨询","source_conversation":1}}]

会话记录：
{conversations_content}
"""


# 会话高价值评估 prompt（用于 AI 智能推荐高价值会话）
CONVERSATION_RECOMMEND_PROMPT_TEMPLATE = """你是客服会话价值评估助手。以下是若干会话的概要信息。请评估每个会话是否包含可提取的高价值客服知识，输出 JSON 数组。

评估标准（满足任一即视为高价值）：
- 包含产品规格、价格、库存等可复用信息
- 包含售后政策、退换货规则、发货说明等可复用信息
- 包含买家常见问题的解答
- 对话深度足够（多轮交互、有实质内容）
- 跳过纯问候、纯表情、纯图片无文字、纯砍价无结论的对话

要求：
1. 每个高价值会话输出：
   - conversation_index：会话序号（来自输入的"会话 N"中的 N）
   - reason：一句话说明推荐理由（不超过 30 字）
   - estimated_value：预估价值评分（1-5，5 为最高）
2. 低价值会话不要输出
3. 最多推荐 20 个高价值会话
4. 输出严格的 JSON 数组格式，无其他文字

输出格式示例：
[{{"conversation_index":1,"reason":"包含退换货政策说明","estimated_value":4}}]

会话概要：
{conversations_summary}
"""


def _authenticated_tenant_id(current_user: dict) -> int:
    try:
        tenant_id = int(current_user.get("tenant_id") or 0)
    except (TypeError, ValueError):
        tenant_id = 0
    if tenant_id <= 0:
        raise ValueError("缺少租户上下文")
    return tenant_id


def _optional_positive_id(value: object, field: str) -> int | None:
    if value in (None, "", 0, "0"):
        return None
    try:
        parsed = int(value)
    except (TypeError, ValueError) as exc:
        raise KnowledgeUploadError(f"{field} 必须为正整数", 400) from exc
    if parsed <= 0:
        raise KnowledgeUploadError(f"{field} 必须为正整数", 400)
    return parsed


def _bind_internal_tenant(header_tenant_id: object, claimed_tenant_id: object) -> int:
    """Use the authenticated internal header as the only tenant authority."""

    tenant_id = _optional_positive_id(header_tenant_id, "内部租户上下文")
    if tenant_id is None:
        raise KnowledgeUploadError("缺少内部租户上下文", 400)
    claimed = _optional_positive_id(claimed_tenant_id, "tenantId")
    if claimed is not None and claimed != tenant_id:
        raise KnowledgeUploadError("请求租户上下文不一致", 403)
    return tenant_id


def _validate_claimed_tenant(claimed_tenant_id: object, authenticated_tenant_id: int) -> None:
    claimed = _optional_positive_id(claimed_tenant_id, "tenantId")
    if claimed is not None and claimed != authenticated_tenant_id:
        raise KnowledgeUploadError("请求租户上下文不一致", 403)


def _validate_declared_media_type(extension: str, declared_media_type: str | None) -> None:
    declared = str(declared_media_type or "").split(";", 1)[0].strip().lower()
    if not declared:
        return
    allowed = _ALLOWED_MEDIA_TYPES.get(extension, set())
    if declared not in allowed:
        raise KnowledgeUploadError("文件媒体类型与扩展名不匹配", 400)


async def _read_upload_limited(file: UploadFile) -> bytes:
    content = await file.read(MAX_FILE_SIZE + 1)
    if not content:
        raise KnowledgeUploadError("文件内容为空", 400)
    if len(content) > MAX_FILE_SIZE:
        raise KnowledgeUploadError("文件不能超过 10MB", 413)
    return bytes(content)


def _validate_ooxml_package(extension: str, content: bytes) -> None:
    """Reject malformed, mismatched, encrypted, or expansion-heavy OOXML."""

    if not content.startswith(b"PK\x03\x04") or not zipfile.is_zipfile(io.BytesIO(content)):
        raise KnowledgeUploadError("Office 文件内容与扩展名不匹配")
    required_member = {
        ".xlsx": "xl/workbook.xml",
        ".pptx": "ppt/presentation.xml",
        ".docx": "word/document.xml",
    }.get(extension)
    if required_member is None:
        raise KnowledgeUploadError("不支持的 Office 文件类型", 400)

    try:
        with zipfile.ZipFile(io.BytesIO(content), "r") as archive:
            members = archive.infolist()
            if len(members) > MAX_OOXML_MEMBERS:
                raise KnowledgeUploadError("Office 文件包含过多内部条目")

            names: set[str] = set()
            total_uncompressed = 0
            for member in members:
                name = member.filename
                if not name or "\\" in name or "\x00" in name:
                    raise KnowledgeUploadError("Office 文件包含不安全路径")
                path = PurePosixPath(name)
                if path.is_absolute() or ".." in path.parts or ":" in path.parts[0]:
                    raise KnowledgeUploadError("Office 文件包含不安全路径")
                if name in names:
                    raise KnowledgeUploadError("Office 文件包含重复条目")
                names.add(name)
                if member.flag_bits & 0x1:
                    raise KnowledgeUploadError("不支持加密的 Office 文件")
                mode = (member.external_attr >> 16) & 0xFFFF
                if mode and stat.S_ISLNK(mode):
                    raise KnowledgeUploadError("Office 文件包含不安全链接")
                if member.is_dir():
                    continue
                if member.file_size < 0 or member.file_size > MAX_OOXML_MEMBER_BYTES:
                    raise KnowledgeUploadError("Office 文件内部条目过大")
                total_uncompressed += member.file_size
                if total_uncompressed > MAX_OOXML_UNCOMPRESSED_BYTES:
                    raise KnowledgeUploadError("Office 文件展开后超过 50MB")
                if member.file_size:
                    ratio = member.file_size / max(member.compress_size, 1)
                    if ratio > MAX_OOXML_COMPRESSION_RATIO:
                        raise KnowledgeUploadError("Office 文件压缩比超过安全限制")
                if name.lower().endswith((".xml", ".rels")):
                    with archive.open(member, "r") as stream:
                        prefix = stream.read(min(member.file_size, 64 * 1024)).upper()
                    if b"<!DOCTYPE" in prefix or b"<!ENTITY" in prefix:
                        raise KnowledgeUploadError("Office 文件包含不安全 XML 声明")

            if "[Content_Types].xml" not in names or required_member not in names:
                raise KnowledgeUploadError("Office 文件结构与扩展名不匹配")
    except KnowledgeUploadError:
        raise
    except (OSError, RuntimeError, zipfile.BadZipFile, zipfile.LargeZipFile) as exc:
        raise KnowledgeUploadError("Office 文件已损坏或无法安全解析") from exc


def _decode_text(content: bytes) -> str:
    if b"\x00" in content:
        raise KnowledgeUploadError("文本文件包含二进制内容")
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            text_value = content.decode(encoding)
            disallowed_controls = sum(
                1 for char in text_value if ord(char) < 32 and char not in "\r\n\t"
            )
            if disallowed_controls > max(4, len(text_value) // 1000):
                raise KnowledgeUploadError("文本文件包含过多控制字符")
            return text_value
        except UnicodeDecodeError:
            continue
    raise KnowledgeUploadError("文本文件编码不受支持，请使用 UTF-8 或 GB18030")


def _release_parse_slot(task: asyncio.Task) -> None:
    try:
        task.exception()
    except (asyncio.CancelledError, Exception):
        pass
    _parse_semaphore.release()


async def _parse_file_async(ext: str, content: bytes, filename: str) -> str:
    await _parse_semaphore.acquire()
    task = asyncio.create_task(asyncio.to_thread(_parse_file, ext, content, filename))
    release_in_finally = True
    try:
        return await asyncio.wait_for(asyncio.shield(task), timeout=PARSE_TIMEOUT_SECONDS)
    except asyncio.TimeoutError as exc:
        release_in_finally = False
        task.add_done_callback(_release_parse_slot)
        raise KnowledgeUploadError("文件解析超时或复杂度超过安全限制") from exc
    except asyncio.CancelledError:
        release_in_finally = False
        task.add_done_callback(_release_parse_slot)
        raise
    finally:
        if release_in_finally:
            _parse_semaphore.release()


async def _prepare_uploaded_file(file: UploadFile) -> tuple[str, str, bytes]:
    filename = str(file.filename or "unknown")
    extension = _get_extension(filename)
    if extension not in ALLOWED_EXTENSIONS:
        raise KnowledgeUploadError(
            f"不支持的文件格式：{extension or '无扩展名'}，仅支持 "
            + "/".join(sorted(ALLOWED_EXTENSIONS)),
            400,
        )
    _validate_declared_media_type(extension, getattr(file, "content_type", None))
    content = await _read_upload_limited(file)
    return filename, extension, content


async def _precheck_knowledge_usage(tenant_id: int, user_id: int, prompt: str) -> None:
    try:
        prompt_tokens = estimate_text_tokens(prompt)
        result = await precheck_ai_usage({
            "tenantId": tenant_id,
            "userId": user_id,
            "scene": "knowledge_base_extract",
            "providerName": "default",
            "modelName": "default",
            "modelType": "chat",
            "billingMode": "token",
            "promptTokens": prompt_tokens,
            "completionTokens": 2_000,
            "totalTokens": prompt_tokens + 2_000,
        })
        if not isinstance(result, dict) or result.get("skipped") is True:
            raise RuntimeError("billing precheck was skipped")
        if result.get("ok") is False or result.get("enough") is not True:
            raise RuntimeError(str(result.get("message") or "billing precheck rejected"))
    except AiBillingError as exc:
        raise KnowledgeUploadError(exc.user_message, exc.status_code) from exc
    except Exception as exc:
        raise KnowledgeUploadError("AI 计费预检查暂时不可用，尚未调用模型，请稍后重试", 503) from exc


async def _charge_knowledge_usage(
    *,
    tenant_id: int,
    user_id: int,
    prompt: str,
    extracted: str,
    ai_result: dict[str, Any],
) -> None:
    provider_request_id = str(ai_result.get("requestId") or "").strip()
    if not provider_request_id:
        raise KnowledgeUploadError("AI 返回结果缺少计费请求编号，结果未交付，请稍后重试", 503)
    try:
        charged = await charge_text_usage(
            tenant_id=tenant_id,
            user_id=user_id,
            scene="knowledge_base_extract",
            provider_name=str(ai_result.get("provider") or "default"),
            model_name=str(ai_result.get("model") or "default"),
            prompt=prompt,
            completion=extracted,
            request_id=provider_request_id,
            raw_usage=ai_result.get("usage") or {},
        )
        if not isinstance(charged, dict) or charged.get("skipped") is True:
            raise RuntimeError("billing charge was skipped")
        if charged.get("deducted") is not True and charged.get("duplicate") is not True:
            raise RuntimeError("billing charge was not confirmed")
    except AiBillingPaymentRequired as exc:
        raise KnowledgeUploadError(
            "AI 已完成处理，但 Token 余额不足；结果未交付，请充值后使用原请求重试",
            402,
        ) from exc
    except AiBillingError as exc:
        raise KnowledgeUploadError(
            "AI 已完成处理，但计费服务暂时不可用；结果未交付，请使用原请求重试",
            exc.status_code,
        ) from exc
    except Exception as exc:
        raise KnowledgeUploadError(
            "AI 已完成处理，但计费结果暂时无法确认；结果未交付，请使用原请求重试",
            503,
        ) from exc


@router.post("/extract")
async def extract_knowledge_base(
    file: UploadFile = File(...),
    userId: Optional[str] = Form(None),
    tenantId: Optional[str] = Form(None),
    db: AsyncSession = Depends(get_db),
    x_internal_tenant_id: Optional[str] = Header(None, alias="X-Internal-Tenant-Id"),
    _: None = Depends(verify_internal_token),
):
    """接收文件，解析后调用 AI 提取客服回复规则，返回结构化 Markdown。

    可选表单字段：
      userId:  当前用户 ID（由 Java 网关从 TenantContext 透传，用于 AI 调用扣费）
      tenantId: 当前租户 ID（同上）
    """
    filename = "unknown"
    try:
        authenticated_tenant_id = _bind_internal_tenant(x_internal_tenant_id, tenantId)
        authenticated_user_id = _optional_positive_id(userId, "userId")
        if authenticated_user_id is None:
            raise KnowledgeUploadError("缺少可计费用户上下文", 400)
        filename, ext, content = await _prepare_uploaded_file(file)

        async with govern_transient_upload(
            content,
            tenant_id=authenticated_tenant_id,
            user_id=authenticated_user_id,
            source_type="knowledge-base-extract",
        ):
            file_text = await _parse_file_async(ext, content, filename)
            if not file_text.strip():
                raise KnowledgeUploadError("文件内容为空或无法解析")
            if len(file_text) > 30000:
                file_text = file_text[:30000] + "\n\n（文件内容过长，已截断）"
                logger.info("知识库文件内容截断至 30000 字符")

            prompt = EXTRACT_PROMPT_TEMPLATE.format(file_content=file_text)
            await _precheck_knowledge_usage(authenticated_tenant_id, authenticated_user_id, prompt)
            extracted, ai_result = await _call_ai_extract(prompt)
            await _charge_knowledge_usage(
                tenant_id=authenticated_tenant_id,
                user_id=authenticated_user_id,
                prompt=prompt,
                extracted=extracted,
                ai_result=ai_result,
            )
            if not extracted or not extracted.strip():
                return ResultObject.failed("AI 未能从文件中提取有效规则，请检查文件内容或重试")

            rule_count = extracted.count("\n- ") + (1 if extracted.strip().startswith("- ") else 0)
            return ResultObject.success({
                "extractedText": extracted.strip(),
                "ruleCount": max(rule_count, 0),
                "fileName": filename,
            })
    except KnowledgeUploadError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except UploadGovernanceError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except Exception as e:
        return safe_route_failure(logger, e, operation="extract knowledge base file", user_message="文件处理失败，请稍后重试")


async def _call_ai_extract(prompt: str) -> tuple[str, Dict[str, Any]]:
    """调用 AI 生成文本。

    ai_provider.generate_text 是 async 函数，返回 dict：
    - 成功：{"ok": True, "content": "...", "provider": ..., "model": ..., "requestId": ..., "usage": ...}
    - 失败：{"ok": False, "error": "...", ...}
    本函数提取 content 文本，失败时抛 RuntimeError 由上层捕获。
    返回元组：(content 文本, 原始 ai_result dict) 供上层扣费使用。
    """
    result: Dict[str, Any] = await generate_text(
        scene="knowledge_base_extract",
        system_prompt="你是专业的客服规则提取助手，擅长从文档中提炼结构化的客服回复规则。",
        user_prompt=prompt,
        temperature=0.3,
    )

    if not result.get("ok"):
        err = result.get("error") or "AI 调用失败且未返回错误信息"
        raise RuntimeError(f"AI 提取失败：{err}")

    content = result.get("content") or ""
    return str(content).strip(), result


# ============================================================
# 新建知识库弹窗 - Q&A 提取接口（文件上传 / 会话聊天提取 / AI 推荐会话）
# ============================================================

import json as _json
from datetime import datetime as _datetime


def _parse_ai_json_array(content: str) -> list[dict]:
    """容错解析 AI 返回的 JSON 数组（去除 markdown 代码块包装）。"""
    if not content:
        return []
    cleaned = content.strip()
    # 去除 markdown 代码块包装
    if cleaned.startswith("```"):
        # 去掉第一行（```json 或 ```）
        cleaned = cleaned.split("\n", 1)[-1] if "\n" in cleaned else cleaned[3:]
        # 去掉末尾 ```
        if cleaned.rstrip().endswith("```"):
            cleaned = cleaned.rstrip()[:-3]
    # 截取第一个 [ 到最后一个 ] 之间的内容
    start = cleaned.find("[")
    end = cleaned.rfind("]")
    if start == -1 or end == -1 or end <= start:
        return []
    json_str = cleaned[start:end + 1]
    try:
        parsed = _json.loads(json_str)
        if isinstance(parsed, list):
            return parsed
    except (ValueError, _json.JSONDecodeError) as exc:
        logger.warning("AI 返回的 JSON 解析失败: %s", exc)
    return []


def _sanitize_qa_entry(entry: dict, source_label: str = "") -> dict:
    """清洗单条 Q&A 条目，确保字段类型与长度合法。"""
    question = str(entry.get("question") or "").strip()
    answer = str(entry.get("answer") or "").strip()
    if not question or not answer:
        return {}
    return {
        "title": question[:100],
        "content": answer[:5000],
        "category": str(entry.get("category") or "").strip()[:64],
        "tags": str(entry.get("tags") or "").strip()[:512],
        "source_summary": str(entry.get("source_summary") or source_label).strip()[:255],
        "source_conversation": entry.get("source_conversation"),
    }


@router.post("/extract-qa-from-file")
async def extract_qa_from_file(
    file: UploadFile = File(...),
    userId: Optional[str] = Form(None),
    tenantId: Optional[str] = Form(None),
    fileType: Optional[str] = Form(None),
    db: AsyncSession = Depends(get_db),
    x_internal_tenant_id: Optional[str] = Header(None, alias="X-Internal-Tenant-Id"),
    _: None = Depends(verify_internal_token),
):
    """新建知识库弹窗 - 文件上传模式：解析文件 + AI 提取 Q&A 对，返回 JSON 数组（不写入数据库）。

    可选表单字段：
      userId:  当前用户 ID（由 Java 网关透传，用于 AI 调用扣费）
      tenantId: 当前租户 ID
      fileType: 文件类型（auto/chat_records/product_docs/company_docs/general）
                - auto（默认）：自动检测文件类型后选用对应 prompt
                - chat_records：聊天记录文件
                - product_docs：商品资料/说明书
                - company_docs：公司资料/政策文档
                - general：通用提取
    """
    filename = "unknown"
    try:
        authenticated_tenant_id = _bind_internal_tenant(x_internal_tenant_id, tenantId)
        authenticated_user_id = _optional_positive_id(userId, "userId")
        if authenticated_user_id is None:
            raise KnowledgeUploadError("缺少可计费用户上下文", 400)
        filename, ext, content = await _prepare_uploaded_file(file)

        # 规范化 fileType 参数
        file_type = (fileType or "auto").strip().lower()
        if file_type not in ("auto", "chat_records", "product_docs", "company_docs", "general"):
            file_type = "auto"

        async with govern_transient_upload(
            content,
            tenant_id=authenticated_tenant_id,
            user_id=authenticated_user_id,
            source_type="knowledge-base-extract-qa",
        ):
            file_text = await _parse_file_async(ext, content, filename)
            if not file_text.strip():
                raise KnowledgeUploadError("文件内容为空或无法解析")
            if len(file_text) > 30000:
                file_text = file_text[:30000] + "\n\n（文件内容过长，已截断）"
                logger.info("Q&A 提取文件内容截断至 30000 字符")

            # 文件类型检测（auto 模式先调用 AI 识别）
            detected_type = file_type
            if file_type == "auto":
                detected_type = await _detect_file_type(
                    file_text, authenticated_tenant_id, authenticated_user_id
                )
                logger.info("文件类型自动检测: filename=%s -> %s", filename, detected_type)

            # 根据文件类型选择提取 prompt
            prompt = _select_extract_prompt(detected_type, file_text)

            await _precheck_knowledge_usage(authenticated_tenant_id, authenticated_user_id, prompt)
            extracted, ai_result = await _call_ai_extract_qa(prompt)
            await _charge_knowledge_usage(
                tenant_id=authenticated_tenant_id,
                user_id=authenticated_user_id,
                prompt=prompt,
                extracted=extracted,
                ai_result=ai_result,
            )

            entries_raw = _parse_ai_json_array(extracted)
            entries = [
                sanitized for entry in entries_raw
                if (sanitized := _sanitize_qa_entry(entry, source_label=filename))
            ]
            return ResultObject.success({
                "entries": entries,
                "totalCount": len(entries),
                "fileName": filename,
                "fileType": EXTENSION_TYPE_LABEL.get(ext, "未知"),
                "contentCategory": detected_type,
                "rawLength": len(file_text),
            })
    except KnowledgeUploadError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except UploadGovernanceError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except Exception as e:
        return safe_route_failure(logger, e, operation="extract qa from file", user_message="文件处理失败，请稍后重试")


def _select_extract_prompt(file_type: str, file_text: str) -> str:
    """根据文件类型选择对应的提取 prompt。"""
    if file_type == "chat_records":
        return CHAT_RECORDS_QA_EXTRACT_PROMPT_TEMPLATE.format(file_content=file_text)
    if file_type == "product_docs":
        return PRODUCT_DOCS_QA_EXTRACT_PROMPT_TEMPLATE.format(file_content=file_text)
    if file_type == "company_docs":
        return COMPANY_DOCS_QA_EXTRACT_PROMPT_TEMPLATE.format(file_content=file_text)
    # general / unknown 都走通用模板
    return QA_EXTRACT_PROMPT_TEMPLATE.format(file_content=file_text)


async def _detect_file_type(file_text: str, tenant_id: int, user_id: int) -> str:
    """通过 AI 自动检测文件类型。失败时回退为 general。"""
    try:
        sample = file_text[:2000]
        detect_prompt = FILE_TYPE_DETECT_PROMPT_TEMPLATE.format(file_content=sample)
        # 检测调用也走计费预检（但 prompt 较短，token 消耗小）
        await _precheck_knowledge_usage(tenant_id, user_id, detect_prompt)
        result: Dict[str, Any] = await generate_text(
            scene="knowledge_base_detect",
            system_prompt="你是文件类型识别助手，只输出类型标识符，不要其他文字。",
            user_prompt=detect_prompt,
            temperature=0.0,
        )
        if result.get("ok"):
            content = (result.get("content") or "").strip().lower()
            # 容忍 AI 返回带 markdown 代码块或多余文字
            for valid_type in ("chat_records", "product_docs", "company_docs", "general"):
                if valid_type in content:
                    return valid_type
        # 扣费（即使识别失败也扣，因为调用了 AI）
        await _charge_knowledge_usage(
            tenant_id=tenant_id,
            user_id=user_id,
            prompt=detect_prompt,
            extracted="",
            ai_result=result,
        )
    except Exception as e:
        logger.warning("文件类型自动检测失败，回退为 general: %s", e)
    return "general"


async def _call_ai_extract_qa(prompt: str) -> tuple[str, Dict[str, Any]]:
    """调用 AI 提取 Q&A JSON 数组。"""
    result: Dict[str, Any] = await generate_text(
        scene="knowledge_base_extract",
        system_prompt="你是专业的客服知识库提取助手，擅长从文档与会话中提炼问答对，并输出严格的 JSON 数组。",
        user_prompt=prompt,
        temperature=0.3,
    )
    if not result.get("ok"):
        err = result.get("error") or "AI 调用失败且未返回错误信息"
        raise RuntimeError(f"AI 提取失败：{err}")
    content = result.get("content") or ""
    return str(content).strip(), result


@router.post("/extract-qa-from-conversations")
async def extract_qa_from_conversations(
    data: dict = {},
    db: AsyncSession = Depends(get_db),
    x_internal_tenant_id: Optional[str] = Header(None, alias="X-Internal-Tenant-Id"),
    _: None = Depends(verify_internal_token),
):
    """新建知识库弹窗 - 会话聊天提取模式：根据 sids 拉取消息 + AI 提取 Q&A 对。

    请求体：
    {
        "userId": 123,
        "tenantId": 1,
        "accountId": 456,
        "sids": ["xxx@goofish", "yyy@goofish"],   # 必填，至少 1 个（推荐方式）
        "conversationIds": [101, 102],             # 兼容旧字段，若 sids 缺失则从 conversations meta 反查 sid
        "conversations": [                          # 可选，前端直接传入会话概要（避免二次拉取）
            {"sid": "xxx@goofish", "peerUserName": "张三", "goodsTitle": "..."}
        ]
    }
    """
    try:
        authenticated_tenant_id = _bind_internal_tenant(x_internal_tenant_id, data.get("tenantId"))
        authenticated_user_id = _optional_positive_id(data.get("userId"), "userId")
        if authenticated_user_id is None:
            raise KnowledgeUploadError("缺少可计费用户上下文", 400)
        account_id = _optional_positive_id(data.get("accountId"), "accountId")
        if account_id is None:
            raise KnowledgeUploadError("accountId 必填", 400)

        # 优先使用 sids（字符串数组，总是可用）；兼容旧 conversationIds（int 数组，可能为 null）
        sids_input = data.get("sids") or []
        conversations_meta = data.get("conversations") or []

        # 建立 sid → meta 索引（sid 总是存在且唯一）
        meta_by_sid: dict[str, dict] = {}
        for c in conversations_meta:
            if not isinstance(c, dict):
                continue
            sid_val = str(c.get("sid") or "").strip()
            if sid_val:
                meta_by_sid[sid_val] = c

        # 如果 sids 为空但 conversationIds 提供，从 meta 中按 conversationId 反查 sid
        if not sids_input:
            legacy_ids = data.get("conversationIds") or []
            if not isinstance(legacy_ids, list) or not legacy_ids:
                raise KnowledgeUploadError("sids 或 conversationIds 必须为非空数组", 400)
            for conv_id in legacy_ids:
                try:
                    conv_id_int = int(conv_id)
                except (TypeError, ValueError):
                    continue
                for meta in conversations_meta:
                    if isinstance(meta, dict) and meta.get("conversationId") == conv_id_int:
                        sid_val = str(meta.get("sid") or "").strip()
                        if sid_val:
                            sids_input.append(sid_val)
                        break

        if not isinstance(sids_input, list) or not sids_input:
            raise KnowledgeUploadError("sids 必须为非空数组", 400)
        if len(sids_input) > 50:
            raise KnowledgeUploadError("单次最多提取 50 个会话", 400)

        # 拉取每个会话的消息
        from ....services.ws_storage import get_context_messages
        conv_blocks: list[str] = []
        conv_labels: dict[int, str] = {}
        total_messages = 0
        max_messages_per_conv = 100  # 单会话最多 100 条消息
        max_total_messages = 1000    # 总消息上限

        for idx, raw_sid in enumerate(sids_input, start=1):
            if total_messages >= max_total_messages:
                logger.info("会话消息总数已达上限 %d，停止拉取", max_total_messages)
                break
            sid = str(raw_sid or "").strip()
            if not sid:
                continue
            # 去掉 sid: 前缀
            if sid.startswith("sid:"):
                sid = sid[4:]
            meta = meta_by_sid.get(sid) or {}
            peer_user_id = str(meta.get("peerUserId") or meta.get("peer_user_id") or "").strip()
            messages, _ = await get_context_messages(
                db, authenticated_tenant_id, account_id, sid,
                limit=max_messages_per_conv, offset=0,
                user_id=None, peer_user_id=peer_user_id or None,
            )
            if not messages:
                continue
            peer_name = str(meta.get("peerUserName") or "").strip()
            goods_title = str(meta.get("goodsTitle") or "").strip()
            label = f"{peer_name or '买家'} - {goods_title}" if goods_title else (peer_name or f"会话{idx}")
            conv_labels[idx] = label

            # 构造对话文本（每条消息一行）
            msg_lines: list[str] = []
            for m in messages:
                direction = str(m.get("direction") or "").upper()
                sender = str(m.get("senderUserName") or "").strip()
                content_text = str(m.get("msgContent") or m.get("completeMsg") or "").strip()
                if not content_text:
                    continue
                role = "[卖家]" if direction == "OUT" else "[买家]"
                msg_lines.append(f"  {role} {sender}: {content_text[:500]}")
            if msg_lines:
                conv_blocks.append(f"对话 {idx}（{label}）：\n" + "\n".join(msg_lines))
                total_messages += len(msg_lines)

        if not conv_blocks:
            return ResultObject.failed("所选会话均无有效消息，无法提取知识")

        conversations_content = "\n\n".join(conv_blocks)
        if len(conversations_content) > 30000:
            conversations_content = conversations_content[:30000] + "\n\n（会话内容过长，已截断）"

        prompt = CONVERSATION_QA_EXTRACT_PROMPT_TEMPLATE.format(conversations_content=conversations_content)
        await _precheck_knowledge_usage(authenticated_tenant_id, authenticated_user_id, prompt)
        extracted, ai_result = await _call_ai_extract_qa(prompt)
        await _charge_knowledge_usage(
            tenant_id=authenticated_tenant_id,
            user_id=authenticated_user_id,
            prompt=prompt,
            extracted=extracted,
            ai_result=ai_result,
        )

        entries_raw = _parse_ai_json_array(extracted)
        entries: list[dict] = []
        for entry in entries_raw:
            source_conv = entry.get("source_conversation")
            label = conv_labels.get(int(source_conv)) if source_conv else ""
            sanitized = _sanitize_qa_entry(entry, source_label=label)
            if sanitized:
                entries.append(sanitized)

        return ResultObject.success({
            "entries": entries,
            "totalCount": len(entries),
            "selectedCount": len(sids_input),
            "processedCount": len(conv_blocks),
            "totalMessages": total_messages,
        })
    except KnowledgeUploadError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except Exception as e:
        return safe_route_failure(logger, e, operation="extract qa from conversations", user_message="会话提取失败，请稍后重试")


@router.post("/recommend-conversations")
async def recommend_conversations(
    data: dict = {},
    db: AsyncSession = Depends(get_db),
    x_internal_tenant_id: Optional[str] = Header(None, alias="X-Internal-Tenant-Id"),
    _: None = Depends(verify_internal_token),
):
    """新建知识库弹窗 - AI 智能推荐高价值会话。

    请求体：
    {
        "userId": 123,
        "tenantId": 1,
        "accountId": 456,
        "conversations": [    # 前端直接传入会话列表（避免后端二次拉取）
            {"sid": "xxx@goofish", "peerUserName": "张三", "goodsTitle": "...",
             "lastMessage": "...", "messageCount": 12, "lastMessageTime": 1690000000000}
        ]
    }

    优化策略（解决大列表 AI 推荐卡顿问题）：
    1. 启发式预过滤：跳过明显低价值会话（消息数<3、最后消息为纯表情/问候）
    2. 分批处理：每批最多 30 个会话，并行调用 AI（避免单次 prompt 过大导致超时）
    3. 结果合并：合并各批次的推荐结果，按 estimatedValue 降序排列，取前 20 个
    """
    try:
        authenticated_tenant_id = _bind_internal_tenant(x_internal_tenant_id, data.get("tenantId"))
        authenticated_user_id = _optional_positive_id(data.get("userId"), "userId")
        if authenticated_user_id is None:
            raise KnowledgeUploadError("缺少可计费用户上下文", 400)
        conversations = data.get("conversations") or []
        if not isinstance(conversations, list) or not conversations:
            return ResultObject.failed("conversations 必须为非空数组")

        # ===== 启发式预过滤：剔除明显低价值会话 =====
        # 规则：消息数 < 2 或最后消息为纯问候/表情的直接跳过
        GREETING_PATTERNS = ("你好", "在吗", "您好", "hi", "hello", "哈喽", "[表情]", "[图片]", "好的", "嗯", "ok")
        filtered_conversations: list[dict] = []
        skipped_count = 0
        for c in conversations:
            if not isinstance(c, dict):
                continue
            msg_count = int(c.get("messageCount") or 0)
            last_msg = str(c.get("lastMessage") or "").strip().lower()
            # 消息数过少直接跳过
            if msg_count < 2:
                skipped_count += 1
                continue
            # 最后消息是纯问候/表情也跳过（但保留消息数多的，可能前面有内容）
            if msg_count < 5 and any(last_msg == pat or last_msg.startswith(pat) for pat in GREETING_PATTERNS):
                skipped_count += 1
                continue
            filtered_conversations.append(c)

        # 如果过滤后为空，回退使用原始列表的前 30 个
        if not filtered_conversations:
            filtered_conversations = conversations[:30]
            logger.info("启发式过滤后为空，回退使用原始前 30 个会话")
        else:
            logger.info(
                "启发式预过滤: 原始 %d 个 → 保留 %d 个（跳过 %d 个低价值）",
                len(conversations), len(filtered_conversations), skipped_count
            )

        # 控制总数上限（避免过多 AI 调用）
        if len(filtered_conversations) > 100:
            filtered_conversations = sorted(
                filtered_conversations,
                key=lambda c: int(c.get("messageCount") or 0),
                reverse=True,
            )[:100]

        # ===== 分批处理：每批 20 个会话（缩小批次以加速单批响应） =====
        # 之前 30 个/批 → 单批 prompt 较长，AI 处理时间约 25-35s
        # 现在 20 个/批 → 单批 prompt 缩短，AI 处理时间约 15-22s
        # 最坏情况：100 个会话 → 5 批 × 22s = 110s（在 180s 超时以内）
        BATCH_SIZE = 20
        batches = [
            filtered_conversations[i:i + BATCH_SIZE]
            for i in range(0, len(filtered_conversations), BATCH_SIZE)
        ]

        # 单批超时控制：避免某批 AI 调用卡住导致整个请求超时
        # 单批上限 50s（含重试），超过则跳过本批继续下一批
        SINGLE_BATCH_TIMEOUT_SECONDS = 50

        async def _process_batch(batch_idx: int, batch: list[dict]) -> list[dict]:
            """处理单个批次，返回 recommendations 列表。"""
            # 使用 1-based idx 作为会话标识（仅在本批次内）
            summary_lines: list[str] = []
            for idx, c in enumerate(batch, start=1):
                peer = str(c.get("peerUserName") or "").strip() or "未知买家"
                goods = str(c.get("goodsTitle") or "").strip() or "未关联商品"
                last_msg = str(c.get("lastMessage") or "").strip()[:80]
                msg_count = int(c.get("messageCount") or 0)
                summary_lines.append(
                    f"会话 {idx} | 买家={peer} | 商品={goods} | 消息数={msg_count} | 最后消息={last_msg}"
                )
            conversations_summary = "\n".join(summary_lines)
            if len(conversations_summary) > 4000:
                conversations_summary = conversations_summary[:4000] + "\n（已截断）"

            prompt = CONVERSATION_RECOMMEND_PROMPT_TEMPLATE.format(conversations_summary=conversations_summary)
            await _precheck_knowledge_usage(authenticated_tenant_id, authenticated_user_id, prompt)
            extracted, ai_result = await _call_ai_extract_qa(prompt)
            await _charge_knowledge_usage(
                tenant_id=authenticated_tenant_id,
                user_id=authenticated_user_id,
                prompt=prompt,
                extracted=extracted,
                ai_result=ai_result,
            )

            batch_recs_raw = _parse_ai_json_array(extracted)
            batch_recs: list[dict] = []
            for rec in batch_recs_raw:
                try:
                    conv_idx = int(rec.get("conversation_index") or rec.get("conversation_id") or 0)
                except (TypeError, ValueError):
                    continue
                if conv_idx <= 0 or conv_idx > len(batch):
                    continue
                meta = batch[conv_idx - 1] if isinstance(batch[conv_idx - 1], dict) else {}
                sid_val = str(meta.get("sid") or "").strip()
                if not sid_val:
                    continue
                batch_recs.append({
                    "sid": sid_val,
                    "conversationId": meta.get("conversationId"),
                    "reason": str(rec.get("reason") or "").strip()[:100],
                    "estimatedValue": max(1, min(5, int(rec.get("estimated_value") or 0))),
                    "peerUserName": meta.get("peerUserName") or "",
                    "goodsTitle": meta.get("goodsTitle") or "",
                    "messageCount": meta.get("messageCount") or 0,
                    "peerUserId": meta.get("peerUserId") or "",
                })
            return batch_recs

        # 串行处理各批次（避免并发 AI 调用触发速率限制）
        # 每批独立超时控制：单批失败/超时不影响后续批次
        all_recommendations: list[dict] = []
        failed_batch_count = 0
        for batch_idx, batch in enumerate(batches):
            try:
                # 使用 asyncio.wait_for 限制单批耗时，避免某批卡住整个请求
                batch_recs = await asyncio.wait_for(
                    _process_batch(batch_idx, batch),
                    timeout=SINGLE_BATCH_TIMEOUT_SECONDS,
                )
                all_recommendations.extend(batch_recs)
                logger.info(
                    "推荐批次 %d/%d 完成，本批 %d 条推荐",
                    batch_idx + 1, len(batches), len(batch_recs)
                )
            except asyncio.TimeoutError:
                failed_batch_count += 1
                logger.warning(
                    "推荐批次 %d/%d 超时（>%ds），跳过继续下一批",
                    batch_idx + 1, len(batches), SINGLE_BATCH_TIMEOUT_SECONDS
                )
                continue
            except Exception as batch_err:
                failed_batch_count += 1
                logger.warning(
                    "推荐批次 %d/%d 失败（跳过）: %s",
                    batch_idx + 1, len(batches), batch_err
                )
                continue

        # 按 estimatedValue 降序排列，取前 20 个
        all_recommendations.sort(key=lambda r: r.get("estimatedValue", 0), reverse=True)
        all_recommendations = all_recommendations[:20]

        # 若所有批次都失败，返回明确的错误提示让前端引导用户重试
        if not all_recommendations and failed_batch_count == len(batches) and len(batches) > 0:
            return ResultObject.failed(
                "AI 推荐服务暂时繁忙，请稍后重试或手动勾选会话",
                code=503,
            )

        return ResultObject.success({
            "recommendations": all_recommendations,
            "totalScanned": len(conversations),
            "filteredCount": len(filtered_conversations),
            "skippedLowValue": skipped_count,
            "batchCount": len(batches),
            "failedBatchCount": failed_batch_count,
            "recommendedCount": len(all_recommendations),
        })
    except KnowledgeUploadError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except Exception as e:
        return safe_route_failure(logger, e, operation="recommend conversations", user_message="AI 推荐失败，请稍后重试")


def _get_extension(filename: str) -> str:
    """获取文件扩展名（小写，带点）。"""
    if "." not in filename:
        return ""
    return "." + filename.rsplit(".", 1)[-1].lower()


def _parse_file(ext: str, content: bytes, filename: str) -> str:
    """根据扩展名选择解析器，提取纯文本。"""
    try:
        if ext in (".md", ".txt"):
            return _decode_text(content)[:MAX_PARSED_CHARACTERS]
        if ext == ".csv":
            return _parse_csv(content)
        if ext == ".xlsx":
            _validate_ooxml_package(ext, content)
            return _parse_excel(content)
        if ext == ".pptx":
            _validate_ooxml_package(ext, content)
            return _parse_ppt(content)
        if ext == ".docx":
            _validate_ooxml_package(ext, content)
            return _parse_docx(content)
        if ext == ".pdf":
            return _parse_pdf(content)
        raise KnowledgeUploadError(f"不支持的文件格式：{ext}", 400)
    except KnowledgeUploadError:
        raise
    except Exception as e:
        logger.warning(
            "知识库文件解析失败 ext=%s errorType=%s",
            ext,
            type(e).__name__,
        )
        raise KnowledgeUploadError("文件已损坏、过于复杂或无法安全解析") from e


def _parse_docx(content: bytes) -> str:
    """解析 Word .docx 文件，返回段落与表格的纯文本。"""
    from docx import Document
    doc = Document(io.BytesIO(content))
    lines: list[str] = []
    total_characters = 0

    # 1. 提取所有段落（按顺序）
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        text = text[:4000]
        remaining = MAX_PARSED_CHARACTERS - total_characters
        if remaining <= 0:
            break
        lines.append(text[:remaining])
        total_characters += min(len(text), remaining) + 1

    # 2. 提取所有表格（每个表格以分隔标记呈现）
    for table_idx, table in enumerate(doc.tables, start=1):
        if total_characters >= MAX_PARSED_CHARACTERS:
            break
        heading = f"### 表格 {table_idx}"
        lines.append(heading)
        total_characters += len(heading) + 1
        for row in table.rows:
            if total_characters >= MAX_PARSED_CHARACTERS:
                break
            cells = [cell.text.strip()[:4000] for cell in row.cells]
            if not any(cells):
                continue
            line = " | ".join(cells)
            remaining = MAX_PARSED_CHARACTERS - total_characters
            if remaining <= 0:
                break
            lines.append(line[:remaining])
            total_characters += min(len(line), remaining) + 1
        lines.append("")

    return "\n".join(lines)


def _parse_pdf(content: bytes) -> str:
    """解析 PDF 文件，返回每页文本（带页码标记）。"""
    import pdfplumber
    lines: list[str] = []
    total_characters = 0
    max_pages = 200  # 安全上限

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        if len(pdf.pages) > max_pages:
            raise KnowledgeUploadError(f"PDF 页数超过安全限制（最多 {max_pages} 页）")
        for page_idx, page in enumerate(pdf.pages, start=1):
            if total_characters >= MAX_PARSED_CHARACTERS:
                break
            try:
                page_text = page.extract_text() or ""
            except Exception as exc:
                logger.warning("PDF 第 %d 页解析失败: %s", page_idx, exc)
                page_text = ""
            if not page_text.strip():
                continue
            heading = f"### 第 {page_idx} 页"
            lines.append(heading)
            total_characters += len(heading) + 1
            # 单页文本截断到 4000 字符避免单页过大
            page_text = page_text[:4000]
            remaining = MAX_PARSED_CHARACTERS - total_characters
            if remaining <= 0:
                break
            lines.append(page_text[:remaining])
            total_characters += min(len(page_text), remaining) + 1
            lines.append("")

    return "\n".join(lines)


def _parse_csv(content: bytes) -> str:
    """解析 CSV 文件，返回表格文本。"""
    text_value = _decode_text(content)
    reader = csv.reader(io.StringIO(text_value))
    lines: list[str] = []
    total_characters = 0
    total_cells = 0
    for row_index, row in enumerate(reader, start=1):
        if row_index > MAX_SPREADSHEET_ROWS:
            raise KnowledgeUploadError("CSV 行数超过安全限制")
        total_cells += len(row)
        if total_cells > MAX_SPREADSHEET_CELLS:
            raise KnowledgeUploadError("CSV 单元格数量超过安全限制")
        if any(cell.strip() for cell in row):
            line = " | ".join(cell[:4000] for cell in row)
            remaining = MAX_PARSED_CHARACTERS - total_characters
            if remaining <= 0:
                break
            lines.append(line[:remaining])
            total_characters += min(len(line), remaining) + 1
    return "\n".join(lines)


def _parse_excel(content: bytes) -> str:
    """解析 Excel 文件，返回所有 sheet 的表格文本。"""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True, keep_links=False)
    try:
        if len(wb.sheetnames) > MAX_WORKSHEETS:
            raise KnowledgeUploadError("工作表数量超过安全限制")
        lines: list[str] = []
        total_rows = 0
        total_cells = 0
        total_characters = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            heading = f"### 工作表：{str(sheet_name)[:200]}"
            lines.append(heading)
            total_characters += len(heading) + 1
            for row in ws.iter_rows(values_only=True):
                total_rows += 1
                total_cells += len(row)
                if total_rows > MAX_SPREADSHEET_ROWS:
                    raise KnowledgeUploadError("工作表总行数超过安全限制")
                if total_cells > MAX_SPREADSHEET_CELLS:
                    raise KnowledgeUploadError("工作表单元格数量超过安全限制")
                values = ["" if cell is None else str(cell)[:4000] for cell in row]
                if not any(value.strip() for value in values):
                    continue
                line = " | ".join(values)
                remaining = MAX_PARSED_CHARACTERS - total_characters
                if remaining <= 0:
                    return "\n".join(lines)
                lines.append(line[:remaining])
                total_characters += min(len(line), remaining) + 1
            lines.append("")
        return "\n".join(lines)
    finally:
        wb.close()


def _parse_ppt(content: bytes) -> str:
    """解析 PPT 文件，返回所有幻灯片文本。"""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    if len(prs.slides) > MAX_SLIDES:
        raise KnowledgeUploadError("幻灯片数量超过安全限制")
    lines: list[str] = []
    total_shapes = 0
    total_table_cells = 0
    total_characters = 0
    for idx, slide in enumerate(prs.slides, start=1):
        total_shapes += len(slide.shapes)
        if total_shapes > MAX_PRESENTATION_SHAPES:
            raise KnowledgeUploadError("幻灯片对象数量超过安全限制")
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        slide_texts.append(paragraph_text[:4000])
            if shape.has_table:
                for row in shape.table.rows:
                    total_table_cells += len(row.cells)
                    if total_table_cells > MAX_PRESENTATION_TABLE_CELLS:
                        raise KnowledgeUploadError("幻灯片表格单元格数量超过安全限制")
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_texts.append(row_text[:8000])
        if slide_texts:
            block = [f"### 幻灯片 {idx}", *slide_texts, ""]
            for line in block:
                remaining = MAX_PARSED_CHARACTERS - total_characters
                if remaining <= 0:
                    return "\n".join(lines)
                lines.append(line[:remaining])
                total_characters += min(len(line), remaining) + 1
    return "\n".join(lines)


# ============================================================
# RAG 知识库接口（基于 SimpleVectorStore 本地向量库）
# ============================================================
@router.post("/rag/add")
async def rag_add(
    data: dict = {},
    current_user: dict = Depends(get_current_user),
):
    """将文本切片向量化并写入 RAG 知识库。

    请求体: {
        "content": "文本内容",
        "goodsId": "123",       # 可选，关联商品 ID
        "source": "manual",      # 可选，来源标识
        "extraMetadata": {...}   # 可选，额外元数据
    }
    """
    try:
        from ....services.rag_service import add_to_rag
        tenant_id = _authenticated_tenant_id(current_user)
        content = data.get("content") or ""
        goods_id = data.get("goodsId")
        source = data.get("source")
        extra_metadata = data.get("extraMetadata")

        if not content.strip():
            return ResultObject.failed("content 不能为空")

        result = await add_to_rag(
            content=content,
            tenant_id=tenant_id,
            goods_id=str(goods_id) if goods_id else None,
            source=source,
            extra_metadata=extra_metadata,
        )
        if result.get("success"):
            return ResultObject.success(result)
        return ResultObject.failed(result.get("error") or "RAG 写入失败")
    except Exception as e:
        if isinstance(e, ValueError) and e.args == ("缺少租户上下文",):
            return ResultObject.failed("缺少租户上下文")
        return safe_route_failure(logger, e, operation="add RAG content", user_message="RAG 写入失败，请稍后重试")


@router.post("/rag/query")
async def rag_query(
    data: dict = {},
    current_user: dict = Depends(get_current_user),
):
    """检索 RAG 知识库中与问题相关的文档片段。

    请求体: {
        "question": "用户问题",
        "goodsId": "123",       # 可选，按商品过滤
        "topK": 5,               # 可选，默认 5
        "similarityThreshold": 0.3  # 可选，默认 0.3
    }
    """
    try:
        from ....services.rag_service import query_rag, DEFAULT_TOP_K, DEFAULT_SIMILARITY_THRESHOLD
        tenant_id = _authenticated_tenant_id(current_user)
        question = data.get("question") or ""
        goods_id = data.get("goodsId")
        top_k = int(data.get("topK") or DEFAULT_TOP_K)
        threshold = float(data.get("similarityThreshold") or DEFAULT_SIMILARITY_THRESHOLD)

        if not question.strip():
            return ResultObject.failed("question 不能为空")

        result = await query_rag(
            question=question,
            tenant_id=tenant_id,
            goods_id=str(goods_id) if goods_id else None,
            top_k=top_k,
            similarity_threshold=threshold,
        )
        if result.get("success"):
            return ResultObject.success(result)
        return ResultObject.failed(result.get("error") or "RAG 检索失败")
    except Exception as e:
        if isinstance(e, ValueError) and e.args == ("缺少租户上下文",):
            return ResultObject.failed("缺少租户上下文")
        return safe_route_failure(logger, e, operation="query RAG content", user_message="RAG 检索失败，请稍后重试")


@router.post("/rag/chat")
async def rag_chat(
    data: dict = {},
    current_user: dict = Depends(get_current_user),
):
    """RAG 完整链路：检索相关文档 + AI 生成回复。

    请求体: {
        "question": "用户问题",
        "goodsId": "123",
        "systemPrompt": "你是客服助手...",  # 可选
        "topK": 5,
        "similarityThreshold": 0.3
    }
    """
    try:
        from ....services.rag_service import chat_by_rag, DEFAULT_TOP_K, DEFAULT_SIMILARITY_THRESHOLD
        tenant_id = _authenticated_tenant_id(current_user)
        question = data.get("question") or ""
        goods_id = data.get("goodsId")
        system_prompt = data.get("systemPrompt")
        top_k = int(data.get("topK") or DEFAULT_TOP_K)
        threshold = float(data.get("similarityThreshold") or DEFAULT_SIMILARITY_THRESHOLD)

        if not question.strip():
            return ResultObject.failed("question 不能为空")

        result = await chat_by_rag(
            question=question,
            tenant_id=tenant_id,
            goods_id=str(goods_id) if goods_id else None,
            system_prompt=system_prompt,
            top_k=top_k,
            similarity_threshold=threshold,
        )
        if result.get("success"):
            return ResultObject.success(result)
        return ResultObject.failed(result.get("error") or "RAG 对话失败")
    except Exception as e:
        if isinstance(e, ValueError) and e.args == ("缺少租户上下文",):
            return ResultObject.failed("缺少租户上下文")
        return safe_route_failure(logger, e, operation="RAG chat", user_message="RAG 对话失败，请稍后重试")


@router.post("/rag/delete")
async def rag_delete(
    data: dict = {},
    current_user: dict = Depends(get_current_user),
):
    """删除指定商品的所有 RAG 文档。

    请求体: {"goodsId": "123"}
    """
    try:
        from ....services.rag_service import delete_rag_by_goods_id
        tenant_id = _authenticated_tenant_id(current_user)
        goods_id = data.get("goodsId")
        if not goods_id:
            return ResultObject.failed("goodsId 不能为空")

        result = await delete_rag_by_goods_id(str(goods_id), tenant_id=tenant_id)
        return ResultObject.success(result)
    except Exception as e:
        if isinstance(e, ValueError) and e.args == ("缺少租户上下文",):
            return ResultObject.failed("缺少租户上下文")
        return safe_route_failure(logger, e, operation="delete RAG content", user_message="RAG 删除失败，请稍后重试")


@router.get("/rag/stats")
async def rag_stats(
    current_user: dict = Depends(get_current_user),
):
    """获取 RAG 知识库统计信息"""
    try:
        from ....services.rag_service import get_rag_stats
        tenant_id = _authenticated_tenant_id(current_user)
        stats = await get_rag_stats(tenant_id=tenant_id)
        return ResultObject.success(stats)
    except Exception as e:
        if isinstance(e, ValueError) and e.args == ("缺少租户上下文",):
            return ResultObject.failed("缺少租户上下文")
        return safe_route_failure(logger, e, operation="get RAG stats", user_message="获取统计失败，请稍后重试")


@router.post("/rag/extract-and-add")
async def rag_extract_and_add(
    file: UploadFile = File(...),
    goodsId: Optional[str] = Form(None),
    userId: Optional[str] = Form(None),
    tenantId: Optional[str] = Form(None),
    current_user: dict = Depends(get_current_user),
):
    """一键完成：上传文件 → AI 提取规则 → 切片向量化 → 写入 RAG

    结合原有的 /extract 接口和 RAG /add 接口，端到端完成知识库入库。
    """
    filename = "unknown"
    try:
        from ....services.rag_service import add_to_rag
        authenticated_tenant_id = _authenticated_tenant_id(current_user)
        authenticated_user_id = _optional_positive_id(current_user.get("user_id"), "userId")
        if authenticated_user_id is None:
            raise KnowledgeUploadError("缺少可计费用户上下文", 400)
        _validate_claimed_tenant(tenantId, authenticated_tenant_id)
        claimed_user_id = _optional_positive_id(userId, "userId")
        if claimed_user_id is not None and claimed_user_id != authenticated_user_id:
            raise KnowledgeUploadError("请求用户上下文不一致", 403)

        filename, ext, content = await _prepare_uploaded_file(file)
        async with govern_transient_upload(
            content,
            tenant_id=authenticated_tenant_id,
            user_id=authenticated_user_id,
            source_type="knowledge-base-rag",
        ):
            file_text = await _parse_file_async(ext, content, filename)
            if not file_text.strip():
                raise KnowledgeUploadError("文件内容为空或无法解析")
            if len(file_text) > 30000:
                file_text = file_text[:30000] + "\n\n（文件内容过长，已截断）"

            prompt = EXTRACT_PROMPT_TEMPLATE.format(file_content=file_text)
            await _precheck_knowledge_usage(authenticated_tenant_id, authenticated_user_id, prompt)
            extracted, ai_result = await _call_ai_extract(prompt)
            await _charge_knowledge_usage(
                tenant_id=authenticated_tenant_id,
                user_id=authenticated_user_id,
                prompt=prompt,
                extracted=extracted,
                ai_result=ai_result,
            )
            if not extracted or not extracted.strip():
                return ResultObject.failed("AI 未能从文件中提取有效规则")

            rag_result = await add_to_rag(
                content=extracted,
                tenant_id=authenticated_tenant_id,
                goods_id=goodsId,
                source=f"file:{filename}",
                extra_metadata={"fileName": filename, "extractedBy": "ai"},
            )

            if not rag_result.get("success"):
                return ResultObject.success({
                    "extractedText": extracted.strip(),
                    "ruleCount": extracted.count("\n- ") + (1 if extracted.strip().startswith("- ") else 0),
                    "fileName": filename,
                    "ragWarning": rag_result.get("error") or "RAG 写入失败",
                })

            return ResultObject.success({
                "extractedText": extracted.strip(),
                "ruleCount": extracted.count("\n- ") + (1 if extracted.strip().startswith("- ") else 0),
                "fileName": filename,
                "ragChunkCount": rag_result.get("chunkCount", 0),
                "ragDocId": rag_result.get("docId"),
                "goodsId": goodsId,
            })
    except KnowledgeUploadError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except UploadGovernanceError as e:
        return ResultObject.failed(e.public_message, code=e.status_code)
    except Exception as e:
        if isinstance(e, ValueError) and e.args == ("缺少租户上下文",):
            return ResultObject.failed("缺少租户上下文")
        return safe_route_failure(logger, e, operation="extract and add RAG content", user_message="处理失败，请稍后重试")
