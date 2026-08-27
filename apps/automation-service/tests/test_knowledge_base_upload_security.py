import asyncio
import io
import time
import zipfile
from contextlib import asynccontextmanager

import pytest

from app.api.v1.routes import knowledge_base


def _ooxml_bytes(*names: str, repeated_payload: bytes | None = None) -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name in names:
            payload = repeated_payload if repeated_payload is not None else b"<root/>"
            archive.writestr(name, payload)
    return output.getvalue()


def test_binary_office_formats_are_not_advertised_as_supported():
    assert ".xls" not in knowledge_base.ALLOWED_EXTENSIONS
    assert ".ppt" not in knowledge_base.ALLOWED_EXTENSIONS
    assert {".md", ".txt", ".csv", ".xlsx", ".pptx", ".docx", ".pdf"} == knowledge_base.ALLOWED_EXTENSIONS


def test_ooxml_preflight_accepts_matching_minimal_packages():
    xlsx = _ooxml_bytes("[Content_Types].xml", "xl/workbook.xml")
    pptx = _ooxml_bytes("[Content_Types].xml", "ppt/presentation.xml")

    knowledge_base._validate_ooxml_package(".xlsx", xlsx)
    knowledge_base._validate_ooxml_package(".pptx", pptx)


def test_pinned_ooxml_parsers_still_accept_real_workbook_and_presentation():
    from openpyxl import Workbook
    from pptx import Presentation

    workbook_bytes = io.BytesIO()
    workbook = Workbook()
    workbook.active.append(["问题", "答案"])
    workbook.active.append(["多久发货", "48 小时内"])
    workbook.save(workbook_bytes)
    workbook.close()

    presentation_bytes = io.BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "售后规则"
    slide.placeholders[1].text = "支持七天退货"
    presentation.save(presentation_bytes)

    assert "48 小时内" in knowledge_base._parse_file(
        ".xlsx", workbook_bytes.getvalue(), "rules.xlsx"
    )
    assert "支持七天退货" in knowledge_base._parse_file(
        ".pptx", presentation_bytes.getvalue(), "rules.pptx"
    )


@pytest.mark.parametrize(
    ("extension", "names"),
    [
        (".xlsx", ("[Content_Types].xml", "ppt/presentation.xml")),
        (".pptx", ("[Content_Types].xml", "xl/workbook.xml")),
        (".xlsx", ("[Content_Types].xml", "../outside.xml", "xl/workbook.xml")),
    ],
)
def test_ooxml_preflight_rejects_wrong_package_type_and_traversal(extension, names):
    with pytest.raises(knowledge_base.KnowledgeUploadError):
        knowledge_base._validate_ooxml_package(extension, _ooxml_bytes(*names))


def test_ooxml_preflight_rejects_excessive_compression_ratio():
    payload = b"A" * (2 * 1024 * 1024)
    archive = _ooxml_bytes("[Content_Types].xml", "xl/workbook.xml", repeated_payload=payload)

    with pytest.raises(knowledge_base.KnowledgeUploadError, match="压缩比"):
        knowledge_base._validate_ooxml_package(".xlsx", archive)


class _Upload:
    filename = "rules.txt"
    content_type = "text/plain"

    def __init__(self, content: bytes):
        self.content = content
        self.read_sizes: list[int] = []

    async def read(self, size: int = -1) -> bytes:
        self.read_sizes.append(size)
        return self.content if size < 0 else self.content[:size]


@pytest.mark.asyncio
async def test_limited_reader_never_requests_more_than_limit_plus_one():
    upload = _Upload(b"A" * (knowledge_base.MAX_FILE_SIZE + 2))

    with pytest.raises(knowledge_base.KnowledgeUploadError, match="10MB"):
        await knowledge_base._read_upload_limited(upload)

    assert upload.read_sizes == [knowledge_base.MAX_FILE_SIZE + 1]


def test_internal_tenant_header_is_authoritative():
    assert knowledge_base._bind_internal_tenant("42", "42") == 42
    assert knowledge_base._bind_internal_tenant("42", None) == 42

    with pytest.raises(knowledge_base.KnowledgeUploadError, match="租户上下文不一致"):
        knowledge_base._bind_internal_tenant("42", "99")
    with pytest.raises(knowledge_base.KnowledgeUploadError, match="租户上下文"):
        knowledge_base._bind_internal_tenant(None, "42")


@pytest.mark.asyncio
async def test_internal_tenant_mismatch_is_rejected_before_reading_file():
    upload = _Upload(b"rules")

    result = await knowledge_base.extract_knowledge_base(
        file=upload,
        userId="7",
        tenantId="99",
        db=None,
        x_internal_tenant_id="42",
        _=None,
    )

    assert result.code == 403
    assert upload.read_sizes == []


@pytest.mark.asyncio
async def test_parser_timeout_keeps_concurrency_slot_until_worker_finishes(monkeypatch):
    monkeypatch.setattr(knowledge_base, "PARSE_TIMEOUT_SECONDS", 0.01)
    monkeypatch.setattr(knowledge_base, "_PARSE_CONCURRENCY", 1)
    monkeypatch.setattr(knowledge_base, "_parse_semaphore", asyncio.Semaphore(1))

    def slow_parse(*_args):
        time.sleep(0.05)
        return "done"

    monkeypatch.setattr(knowledge_base, "_parse_file", slow_parse)
    with pytest.raises(knowledge_base.KnowledgeUploadError, match="超时"):
        await knowledge_base._parse_file_async(".txt", b"hello", "rules.txt")

    assert knowledge_base._parse_semaphore.locked()
    await asyncio.sleep(0.07)
    assert not knowledge_base._parse_semaphore.locked()


@pytest.mark.asyncio
async def test_internal_extract_uses_header_tenant_for_transient_governance(monkeypatch):
    captured = {}

    @asynccontextmanager
    async def governed(content, **kwargs):
        captured.update(kwargs)
        captured["content"] = content
        yield 1

    async def parsed(*_args):
        return "退款规则"

    async def extracted(_prompt):
        return "- 支持七天退货", {
            "provider": "test", "model": "test", "requestId": "provider-request-1", "usage": {}
        }

    async def prechecked(*_args):
        captured["prechecked"] = True

    async def charged(**kwargs):
        captured["chargedTenant"] = kwargs["tenant_id"]
        return {"deducted": True, "requestId": kwargs["request_id"]}

    monkeypatch.setattr(knowledge_base, "govern_transient_upload", governed)
    monkeypatch.setattr(knowledge_base, "_parse_file_async", parsed)
    monkeypatch.setattr(knowledge_base, "_call_ai_extract", extracted)
    monkeypatch.setattr(knowledge_base, "_precheck_knowledge_usage", prechecked)
    monkeypatch.setattr(knowledge_base, "charge_text_usage", charged)

    result = await knowledge_base.extract_knowledge_base(
        file=_Upload(b"refund rules"),
        userId="7",
        tenantId="42",
        db=None,
        x_internal_tenant_id="42",
        _=None,
    )

    assert result.code == 200
    assert captured["tenant_id"] == 42
    assert captured["user_id"] == 7
    assert captured["chargedTenant"] == 42
    assert captured["prechecked"] is True


@pytest.mark.asyncio
async def test_precheck_failure_stops_provider_call(monkeypatch):
    provider_called = False

    @asynccontextmanager
    async def governed(*_args, **_kwargs):
        yield 1

    async def parsed(*_args):
        return "退款规则"

    async def prechecked(*_args):
        raise knowledge_base.KnowledgeUploadError("AI 计费预检查暂时不可用", 503)

    async def extracted(_prompt):
        nonlocal provider_called
        provider_called = True
        return "never", {}

    monkeypatch.setattr(knowledge_base, "govern_transient_upload", governed)
    monkeypatch.setattr(knowledge_base, "_parse_file_async", parsed)
    monkeypatch.setattr(knowledge_base, "_precheck_knowledge_usage", prechecked)
    monkeypatch.setattr(knowledge_base, "_call_ai_extract", extracted)

    result = await knowledge_base.extract_knowledge_base(
        file=_Upload(b"refund rules"),
        userId="7",
        tenantId="42",
        db=None,
        x_internal_tenant_id="42",
        _=None,
    )

    assert result.code == 503
    assert provider_called is False


@pytest.mark.asyncio
async def test_precheck_requires_explicit_billing_approval(monkeypatch):
    async def ambiguous(_payload):
        return {}

    monkeypatch.setattr(knowledge_base, "precheck_ai_usage", ambiguous)
    with pytest.raises(knowledge_base.KnowledgeUploadError, match="尚未调用模型"):
        await knowledge_base._precheck_knowledge_usage(42, 7, "prompt")


@pytest.mark.asyncio
async def test_precheck_preserves_insufficient_balance_state(monkeypatch):
    async def insufficient(_payload):
        raise knowledge_base.AiBillingPaymentRequired("insufficient")

    monkeypatch.setattr(knowledge_base, "precheck_ai_usage", insufficient)
    with pytest.raises(knowledge_base.KnowledgeUploadError) as captured:
        await knowledge_base._precheck_knowledge_usage(42, 7, "prompt")

    assert captured.value.status_code == 402
    assert "Token" in captured.value.public_message


@pytest.mark.asyncio
async def test_charge_requires_confirmation_and_reuses_provider_request_id(monkeypatch):
    captured = {}

    async def unconfirmed(**kwargs):
        captured.update(kwargs)
        return {}

    monkeypatch.setattr(knowledge_base, "charge_text_usage", unconfirmed)
    with pytest.raises(knowledge_base.KnowledgeUploadError, match="计费结果暂时无法确认"):
        await knowledge_base._charge_knowledge_usage(
            tenant_id=42,
            user_id=7,
            prompt="prompt",
            extracted="answer",
            ai_result={"requestId": "provider-request-stable", "usage": {}},
        )

    assert captured["request_id"] == "provider-request-stable"


@pytest.mark.asyncio
async def test_charge_preserves_insufficient_balance_after_provider_call(monkeypatch):
    async def insufficient(**_kwargs):
        raise knowledge_base.AiBillingPaymentRequired("insufficient")

    monkeypatch.setattr(knowledge_base, "charge_text_usage", insufficient)
    with pytest.raises(knowledge_base.KnowledgeUploadError) as captured:
        await knowledge_base._charge_knowledge_usage(
            tenant_id=42,
            user_id=7,
            prompt="prompt",
            extracted="answer",
            ai_result={"requestId": "provider-request-stable", "usage": {}},
        )

    assert captured.value.status_code == 402
    assert "余额不足" in captured.value.public_message


@pytest.mark.asyncio
async def test_rag_extract_does_not_write_or_return_result_when_charge_is_unconfirmed(monkeypatch):
    from app.services import rag_service

    rag_called = False
    precheck_called = False

    @asynccontextmanager
    async def governed(*_args, **_kwargs):
        yield 1

    async def parsed(*_args):
        return "退款规则"

    async def prechecked(*_args):
        nonlocal precheck_called
        precheck_called = True

    async def extracted(_prompt):
        return "- 支持七天退货", {"requestId": "provider-request-2", "usage": {}}

    async def unconfirmed_charge(**_kwargs):
        raise knowledge_base.KnowledgeUploadError("计费结果暂时无法确认", 503)

    async def add_to_rag(**_kwargs):
        nonlocal rag_called
        rag_called = True
        return {"success": True}

    monkeypatch.setattr(knowledge_base, "govern_transient_upload", governed)
    monkeypatch.setattr(knowledge_base, "_parse_file_async", parsed)
    monkeypatch.setattr(knowledge_base, "_precheck_knowledge_usage", prechecked)
    monkeypatch.setattr(knowledge_base, "_call_ai_extract", extracted)
    monkeypatch.setattr(knowledge_base, "_charge_knowledge_usage", unconfirmed_charge)
    monkeypatch.setattr(rag_service, "add_to_rag", add_to_rag)

    result = await knowledge_base.rag_extract_and_add(
        file=_Upload(b"refund rules"),
        goodsId=None,
        userId="7",
        tenantId="42",
        current_user={"user_id": 7, "tenant_id": 42},
    )

    assert result.code == 503
    assert rag_called is False
    assert precheck_called is True
